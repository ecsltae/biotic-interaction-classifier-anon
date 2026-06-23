#!/usr/bin/env python3
"""
Add publication-ready figures to the PowerPoint presentation
"""
from pptx import Presentation
from pptx.util import Inches, Pt
from pathlib import Path

# Load existing presentation
pptx_path = '/path/to/MetaP/classifier/Biotic_Interaction_Classifier_Presentation.pptx'
prs = Presentation(pptx_path)

figures_dir = Path('/path/to/MetaP/classifier/figures')

print("="*70)
print("ADDING FIGURES TO POWERPOINT")
print("="*70)

# Helper function to add image slide
def add_image_slide(title, image_path, description=""):
    slide = prs.slides.add_slide(prs.slide_layouts[5])  # Blank layout

    # Add title
    title_shape = slide.shapes.title
    title_shape.text = title
    title_shape.text_frame.paragraphs[0].font.size = Pt(32)
    title_shape.text_frame.paragraphs[0].font.bold = True

    # Add image - centered and large
    left = Inches(0.5)
    top = Inches(1.5)
    width = Inches(9)

    pic = slide.shapes.add_picture(str(image_path), left, top, width=width)

    # Add description if provided
    if description:
        text_box = slide.shapes.add_textbox(Inches(0.5), Inches(6.8), Inches(9), Inches(0.5))
        tf = text_box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = description
        p.font.size = Pt(12)
        p.font.italic = True

    return slide

# Insert figures after slide 5 (CV Results)
# We'll add them at the end and then they can be reordered if needed

print("\nAdding figure slides...")

# Figure 1: CV Results Comparison
if (figures_dir / 'cv_results_comparison.png').exists():
    add_image_slide(
        "Cross-Validation Performance",
        figures_dir / 'cv_results_comparison.png',
        "BiomedBERT shows the best overall performance across all metrics on 20,000 training samples"
    )
    print("  ✓ Added: CV Results Comparison")

# Figure 2: Confusion Matrix
if (figures_dir / 'confusion_matrix_ensemble.png').exists():
    add_image_slide(
        "Ensemble Model - Confusion Matrix",
        figures_dir / 'confusion_matrix_ensemble.png',
        "Performance on 100-sentence real-world test set: 77% accuracy, 50% precision, 43.5% recall"
    )
    print("  ✓ Added: Confusion Matrix")

# Figure 3: Model Comparison
if (figures_dir / 'model_comparison.png').exists():
    add_image_slide(
        "Test Set Performance Comparison",
        figures_dir / 'model_comparison.png',
        "Ensemble (F1-optimized) provides the best balance of precision and recall"
    )
    print("  ✓ Added: Model Comparison")

# Figure 4: Error Distribution
if (figures_dir / 'error_distribution.png').exists():
    add_image_slide(
        "Error Analysis",
        figures_dir / 'error_distribution.png',
        "77 correct predictions, 10 false positives, 13 false negatives out of 100 test sentences"
    )
    print("  ✓ Added: Error Distribution")

# Figure 5: Probability Distribution
if (figures_dir / 'probability_distribution.png').exists():
    add_image_slide(
        "Prediction Confidence Distribution",
        figures_dir / 'probability_distribution.png',
        "Model confidence scores for true positives vs true negatives, showing threshold placement"
    )
    print("  ✓ Added: Probability Distribution")

# Figure 6: Model Agreement
if (figures_dir / 'model_agreement.png').exists():
    add_image_slide(
        "Model Agreement Analysis",
        figures_dir / 'model_agreement.png',
        "Comparison of BiomedBERT and RoBERTa predictions shows complementary strengths"
    )
    print("  ✓ Added: Model Agreement")

# Save updated presentation
output_path = '/path/to/MetaP/classifier/Biotic_Interaction_Classifier_WITH_FIGURES.pptx'
prs.save(output_path)

print("\n" + "="*70)
print("✓ PRESENTATION UPDATED WITH FIGURES")
print("="*70)
print(f"\nSaved to: {output_path}")
print(f"\nTotal slides: {len(prs.slides)}")
print("\nNew figure slides added:")
print("  • Cross-Validation Performance")
print("  • Ensemble Model - Confusion Matrix")
print("  • Test Set Performance Comparison")
print("  • Error Analysis")
print("  • Prediction Confidence Distribution")
print("  • Model Agreement Analysis")
print("\nYou can reorder these slides in PowerPoint as needed!")